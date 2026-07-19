"""Generate a test .pptx with some slides marked as hidden.

Usage:
    .venv\\Scripts\\python tests\\make_hidden_deck.py

Creates sample/hidden-slides.pptx with 5 slides; slides 2 and 4 are hidden.
Uses python-pptx — Microsoft Office is NOT required to generate the file
(Office is still needed to convert it through doc2md).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

OUT = Path(__file__).resolve().parent.parent / "sample" / "hidden-slides.pptx"
HIDDEN = (2, 4)
TOTAL = 5


def main() -> None:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(1, TOTAL + 1):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.text_frame.text = f"Slide {i}"
    for idx in HIDDEN:
        # A hidden slide carries show="0" on its <p:sld> element.
        prs.slides[idx - 1]._element.set("show", "0")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT} ({TOTAL} slides; hidden: {HIDDEN})")


if __name__ == "__main__":
    main()
